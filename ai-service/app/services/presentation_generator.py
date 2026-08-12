"""Generate PowerPoint presentations from a LessonPlan."""

import os
from pathlib import Path
from uuid import uuid4
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from app.schemas.generation import LessonPlan


class PresentationGenerator:
    """Generate .pptx files based on lesson plans using python-pptx."""

    def __init__(self, upload_dir: str | Path = "../../uploads") -> None:
        self.upload_dir = Path(upload_dir).resolve()
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    def generate(self, lesson_plan: LessonPlan, course_title: str) -> dict[str, str]:
        """
        Generate a PowerPoint file for the lesson and return its metadata.
        Returns a dictionary with 'file_name' and 'file_path' (URL route).
        """
        prs = Presentation()

        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        if title:
            title.text = lesson_plan.presentation_title
        if subtitle:
            subtitle.text = course_title

        # Slide 2: Learning Objectives
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        if title_shape:
            title_shape.text = "Learning Objectives"
        if body_shape:
            tf = body_shape.text_frame
            for objective in lesson_plan.learning_objectives:
                p = tf.add_paragraph()
                p.text = objective
                p.level = 0

        # Slide 3: Content Slide (Generated from description)
        # We will split the description into sentences and distribute them over 2-4 slides
        sentences = [s.strip() for s in lesson_plan.description.replace("\n", " ").split(". ") if s.strip()]
        
        # Group sentences into chunks of 3 for slides
        chunk_size = 3
        for i in range(0, len(sentences), chunk_size):
            chunk = sentences[i:i + chunk_size]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            if title_shape:
                title_shape.text = lesson_plan.title
            if body_shape:
                tf = body_shape.text_frame
                for sentence in chunk:
                    p = tf.add_paragraph()
                    p.text = sentence
                    if not p.text.endswith("."):
                        p.text += "."
                    p.level = 0

        # Slide 4: Conclusion
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        if title:
            title.text = "Summary"
        if subtitle:
            subtitle.text = "End of Lesson"

        # Save the presentation
        file_id = str(uuid4())
        file_name = f"{lesson_plan.title.replace(' ', '_')}_{file_id[:8]}.pptx"
        
        # Remove special characters from file name for safety
        file_name = "".join(c for c in file_name if c.isalnum() or c in ("_", "-", "."))
        file_path = self.upload_dir / file_name

        prs.save(str(file_path))

        # Return the data expected by Next.js /uploads/route.ts
        return {
            "fileName": file_name,
            "filePath": f"/uploads/{file_name}",
        }
