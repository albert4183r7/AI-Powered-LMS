"""Extract images from reference documents for the visual catalog.

Extracted images are filtered by a configurable minimum dimension to
exclude tiny icons and decorative elements that would not be useful
in generated presentations.
"""

import logging
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from uuid import uuid4

LOGGER = logging.getLogger(__name__)

DEFAULT_MIN_IMAGE_DIMENSION = 100


@dataclass(frozen=True)
class ExtractedImage:
    """One image extracted from a reference with provenance metadata."""

    image_id: str
    source_page: int
    width: int
    height: int
    image_format: str
    image_bytes: bytes


def _is_large_enough(width: int, height: int, min_dimension: int) -> bool:
    """Skip images smaller than the configured threshold."""

    return width >= min_dimension and height >= min_dimension


def extract_images_from_pdf(
    content: bytes,
    *,
    min_dimension: int = DEFAULT_MIN_IMAGE_DIMENSION,
) -> list[ExtractedImage]:
    """Extract images from each PDF page using PyMuPDF."""

    import fitz  # PyMuPDF

    extracted_images: list[ExtractedImage] = []
    with fitz.open(stream=content, filetype="pdf") as pdf_document:
        for page_index in range(len(pdf_document)):
            page = pdf_document[page_index]
            for image_info in page.get_images(full=True):
                xref = image_info[0]
                try:
                    base_image = pdf_document.extract_image(xref)
                except Exception:
                    LOGGER.debug("Skipping unextractable PDF image xref=%d.", xref)
                    continue

                width = base_image.get("width", 0)
                height = base_image.get("height", 0)
                image_bytes = base_image.get("image", b"")
                image_ext = base_image.get("ext", "png")

                if not image_bytes or not _is_large_enough(width, height, min_dimension):
                    continue

                extracted_images.append(
                    ExtractedImage(
                        image_id=str(uuid4()),
                        source_page=page_index + 1,
                        width=width,
                        height=height,
                        image_format=image_ext,
                        image_bytes=image_bytes,
                    )
                )

    return extracted_images


def extract_images_from_pptx(
    content: bytes,
    *,
    min_dimension: int = DEFAULT_MIN_IMAGE_DIMENSION,
) -> list[ExtractedImage]:
    """Extract images from PPTX slide shapes."""

    from pptx import Presentation
    from pptx.shapes.picture import Picture

    presentation = Presentation(BytesIO(content))
    extracted_images: list[ExtractedImage] = []

    for slide_index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if not isinstance(shape, Picture):
                continue

            image = shape.image
            image_bytes = image.blob
            content_type = image.content_type or ""

            # Determine format from content type.
            format_map: dict[str, str] = {
                "image/png": "png",
                "image/jpeg": "jpeg",
                "image/gif": "gif",
                "image/bmp": "bmp",
                "image/tiff": "tiff",
                "image/x-emf": "emf",
                "image/x-wmf": "wmf",
            }
            image_format = format_map.get(content_type, "png")

            # Use shape dimensions in EMU (English Metric Units), convert to approximate pixels.
            width_px = int((shape.width or 0) / 9525) if shape.width else 0
            height_px = int((shape.height or 0) / 9525) if shape.height else 0

            if not image_bytes or not _is_large_enough(width_px, height_px, min_dimension):
                continue

            extracted_images.append(
                ExtractedImage(
                    image_id=str(uuid4()),
                    source_page=slide_index + 1,
                    width=width_px,
                    height=height_px,
                    image_format=image_format,
                    image_bytes=image_bytes,
                )
            )

    return extracted_images


def extract_images_from_docx(
    content: bytes,
    *,
    min_dimension: int = DEFAULT_MIN_IMAGE_DIMENSION,
) -> list[ExtractedImage]:
    """Extract inline images from a DOCX file.

    DOCX inline images do not have reliable dimension metadata in the
    relationship, so we attempt to read actual image dimensions from
    the binary content using basic header parsing.
    """

    from docx import Document

    document = Document(BytesIO(content))
    extracted_images: list[ExtractedImage] = []

    for rel_id, rel in document.part.rels.items():
        if "image" not in rel.reltype:
            continue
        try:
            image_bytes = rel.target_part.blob
        except Exception:
            LOGGER.debug("Skipping unreadable DOCX image rel=%s.", rel_id)
            continue

        if not image_bytes:
            continue

        width, height = _guess_image_dimensions(image_bytes)
        if not _is_large_enough(width, height, min_dimension):
            continue

        content_type = rel.target_part.content_type or ""
        image_format = content_type.split("/")[-1] if "/" in content_type else "png"

        extracted_images.append(
            ExtractedImage(
                image_id=str(uuid4()),
                source_page=1,
                width=width,
                height=height,
                image_format=image_format,
                image_bytes=image_bytes,
            )
        )

    return extracted_images


def _guess_image_dimensions(image_bytes: bytes) -> tuple[int, int]:
    """Read image dimensions from binary headers without a full decode.

    Returns (0, 0) when the format is unrecognized, which causes the
    image to be filtered out by the minimum dimension check.
    """

    # PNG: width and height are at bytes 16-23.
    if image_bytes[:8] == b"\x89PNG\r\n\x1a\n" and len(image_bytes) >= 24:
        width = int.from_bytes(image_bytes[16:20], "big")
        height = int.from_bytes(image_bytes[20:24], "big")
        return width, height

    # JPEG: scan for SOF0 marker (0xFF 0xC0).
    if image_bytes[:2] == b"\xff\xd8":
        offset = 2
        while offset < len(image_bytes) - 9:
            if image_bytes[offset] != 0xFF:
                break
            marker = image_bytes[offset + 1]
            if marker == 0xC0 or marker == 0xC2:
                height = int.from_bytes(image_bytes[offset + 5 : offset + 7], "big")
                width = int.from_bytes(image_bytes[offset + 7 : offset + 9], "big")
                return width, height
            segment_length = int.from_bytes(image_bytes[offset + 2 : offset + 4], "big")
            offset += 2 + segment_length
        return 0, 0

    return 0, 0


def save_extracted_images(
    extracted_images: list[ExtractedImage],
    storage_directory: Path,
) -> dict[str, str]:
    """Write extracted images to disk and return a mapping of image_id to file path."""

    storage_directory.mkdir(parents=True, exist_ok=True)
    image_paths: dict[str, str] = {}

    for image in extracted_images:
        image_filename = f"{image.image_id}.{image.image_format}"
        image_path = storage_directory / image_filename
        image_path.write_bytes(image.image_bytes)
        image_paths[image.image_id] = str(image_path)

    return image_paths


def extract_images(
    content: bytes,
    extension: str,
    *,
    min_dimension: int = DEFAULT_MIN_IMAGE_DIMENSION,
) -> list[ExtractedImage]:
    """Dispatch to the appropriate image extractor by file extension."""

    extractor_by_extension = {
        ".pdf": extract_images_from_pdf,
        ".pptx": extract_images_from_pptx,
        ".docx": extract_images_from_docx,
    }

    extractor = extractor_by_extension.get(extension)
    if extractor is None:
        return []

    return extractor(content, min_dimension=min_dimension)
