"""Extract text content from supported reference file formats.

Each extractor preserves source location (page or slide number) so
that Phase 9's RAG pipeline can maintain citation provenance.
"""

from dataclasses import dataclass
from io import BytesIO


@dataclass(frozen=True)
class TextChunk:
    """One block of extracted text tied to its source location."""

    content: str
    source_page: int
    source_location: str


def extract_text_from_pdf(content: bytes) -> list[TextChunk]:
    """Extract text page-by-page using PyMuPDF.

    Each page becomes a separate chunk so citations can reference
    specific pages.
    """

    import fitz  # PyMuPDF

    text_chunks: list[TextChunk] = []
    with fitz.open(stream=content, filetype="pdf") as pdf_document:
        for page_index in range(len(pdf_document)):
            page = pdf_document[page_index]
            page_text = page.get_text("text").strip()
            if page_text:
                text_chunks.append(
                    TextChunk(
                        content=page_text,
                        source_page=page_index + 1,
                        source_location=f"Page {page_index + 1}",
                    )
                )
    return text_chunks


def extract_text_from_docx(content: bytes) -> list[TextChunk]:
    """Extract paragraph text from a DOCX file.

    DOCX files do not have a natural page concept, so all paragraphs
    are grouped into a single chunk with page 1 for simplicity.
    """

    from docx import Document

    document = Document(BytesIO(content))
    paragraphs: list[str] = []
    for paragraph in document.paragraphs:
        paragraph_text = paragraph.text.strip()
        if paragraph_text:
            paragraphs.append(paragraph_text)

    if not paragraphs:
        return []

    return [
        TextChunk(
            content="\n\n".join(paragraphs),
            source_page=1,
            source_location="Document body",
        )
    ]


def extract_text_from_pptx(content: bytes) -> list[TextChunk]:
    """Extract text from each slide's text frames.

    Each slide becomes a separate chunk with its slide number preserved
    for citation.
    """

    from pptx import Presentation

    presentation = Presentation(BytesIO(content))
    text_chunks: list[TextChunk] = []

    for slide_index, slide in enumerate(presentation.slides):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        slide_texts.append(paragraph_text)

        if slide_texts:
            text_chunks.append(
                TextChunk(
                    content="\n".join(slide_texts),
                    source_page=slide_index + 1,
                    source_location=f"Slide {slide_index + 1}",
                )
            )

    return text_chunks


def extract_text_from_txt(content: bytes) -> list[TextChunk]:
    """Read a plain text file with UTF-8 encoding.

    Falls back to latin-1 if UTF-8 decoding fails, since latin-1
    accepts any byte sequence.
    """

    try:
        decoded_text = content.decode("utf-8").strip()
    except UnicodeDecodeError:
        decoded_text = content.decode("latin-1").strip()

    if not decoded_text:
        return []

    return [
        TextChunk(
            content=decoded_text,
            source_page=1,
            source_location="Full text",
        )
    ]


EXTENSION_TEXT_EXTRACTORS: dict[str, type[object] | None] = {
    ".pdf": None,
    ".docx": None,
    ".pptx": None,
    ".txt": None,
}


def extract_text(content: bytes, extension: str) -> list[TextChunk]:
    """Dispatch to the appropriate text extractor by file extension."""

    extractor_by_extension = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx,
        ".txt": extract_text_from_txt,
    }

    extractor = extractor_by_extension.get(extension)
    if extractor is None:
        return []

    return extractor(content)
