#!/usr/bin/env python3
"""
validate_pptx.py — cheap structural checks that catch the most common
generation defects BEFORE you spend a render/vision-QA cycle on them.

This intentionally does not try to catch visual issues (overflow, overlap,
contrast) — that needs actual rendering + a vision model, see
render_preview.sh and your visual-QA step. This script catches things a
generated file should never have: it should just fail fast and cheap.

Usage:
    python3 validate_pptx.py deck.pptx [--expected-slides N]
"""
import argparse
import sys
from pptx import Presentation
from pptx.util import Emu

SLIDE_W = 13.333 * 914400  # EMU for widescreen
SLIDE_H = 7.5 * 914400


def validate(path, expected_slides=None):
    errors = []
    warnings = []

    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"FATAL: could not open {path}: {e}")
        sys.exit(1)

    n_slides = len(prs.slides)
    if expected_slides and n_slides != expected_slides:
        errors.append(f"Expected {expected_slides} slides, found {n_slides}")

    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # Off-slide shapes: usually a sign of a positioning bug — BUT
            # decorative "bleed" shapes (e.g. large background circles that
            # intentionally extend past the edge) are a legitimate pattern.
            # Use a generous tolerance so only clearly-wrong positioning
            # gets flagged; tighten OFFSLIDE_TOLERANCE_IN if you don't use
            # bleed shapes in your component library.
            OFFSLIDE_TOLERANCE_IN = 2.5
            tol = Emu(int(OFFSLIDE_TOLERANCE_IN * 914400))
            if shape.left is not None and shape.top is not None:
                if shape.left < -tol or shape.top < -tol:
                    warnings.append(f"Slide {i}: shape '{shape.shape_type}' positioned far off-slide ({shape.left}, {shape.top})")
                if shape.left is not None and shape.width is not None:
                    if shape.left + shape.width > SLIDE_W + tol:
                        warnings.append(f"Slide {i}: shape extends well past right edge")

            # Empty text frames with a placeholder-looking name often mean
            # a template slot didn't get filled.
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text.lower() in {"lorem ipsum", "todo", "[insert]", "xxx", "placeholder"}:
                    errors.append(f"Slide {i}: leftover placeholder text: '{text}'")

    if errors:
        print(f"FAILED: {len(errors)} error(s)")
        for e in errors:
            print(f"  ERROR   {e}")
        for w in warnings:
            print(f"  WARNING {w}")
        sys.exit(1)

    if warnings:
        print(f"PASSED with {len(warnings)} warning(s)")
        for w in warnings:
            print(f"  WARNING {w}")
    else:
        print(f"PASSED — {n_slides} slides, no issues found")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx_path")
    parser.add_argument("--expected-slides", type=int, default=None)
    args = parser.parse_args()
    validate(args.pptx_path, args.expected_slides)
